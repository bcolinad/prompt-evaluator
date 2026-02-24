"""Document loaders — convert file bytes to text using LangChain document loaders."""

from __future__ import annotations

import asyncio
import csv
import logging
from pathlib import Path
from typing import TYPE_CHECKING

from src.config import get_settings
from src.documents.exceptions import DocumentProcessingError, UnsupportedFormatError
from src.documents.models import DocumentMetadata

if TYPE_CHECKING:
    from langchain_core.documents import Document as LCDocument

logger = logging.getLogger(__name__)

# Supported document extensions and their loader types
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".pdf",
        ".docx",
        ".xlsx",
        ".pptx",
        ".csv",
    }
)


def _count_words(text: str) -> int:
    """Count words in a text string."""
    return len(text.split())


def _pdfplumber_available() -> bool:
    """Check whether pdfplumber is installed."""
    try:
        import pdfplumber  # noqa: F401

        return True
    except ImportError:
        return False


def _pymupdf_available() -> bool:
    """Check whether PyMuPDF (fitz) is installed."""
    try:
        import fitz  # noqa: F401

        return True
    except ImportError:
        return False


def _extract_with_pdfplumber_sync(file_path: Path) -> str:
    """Extract text from a PDF using pdfplumber (sync, called via asyncio.to_thread)."""
    import pdfplumber

    text_parts: list[str] = []
    with pdfplumber.open(str(file_path)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)
    return "\f\n\n".join(text_parts)


def _extract_with_pymupdf_ocr_sync(file_path: Path) -> str:
    """Extract text from a PDF using PyMuPDF OCR (sync, called via asyncio.to_thread)."""
    import fitz

    text_parts: list[str] = []
    doc = fitz.open(str(file_path))
    for page in doc:
        page_text = page.get_text("text") or ""
        if not page_text.strip():
            page_text = page.get_textpage_ocr().extractText() or ""
        if page_text.strip():
            text_parts.append(page_text)
    doc.close()
    return "\f\n\n".join(text_parts)


async def load_document(file_path: Path, filename: str | None = None) -> tuple[str, DocumentMetadata]:
    """Load a document file and extract its text content.

    Uses format-specific LangChain loaders to extract text from various
    document formats. Falls back to plain text reading for unrecognized formats.

    Args:
        file_path: Path to the document file on disk.
        filename: Original filename (for display and extension detection).

    Returns:
        Tuple of (extracted_text, metadata).

    Raises:
        UnsupportedFormatError: If the file extension is not supported.
        DocumentProcessingError: If loading fails.
    """
    if filename is None:
        filename = file_path.name

    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFormatError(suffix, filename=filename)

    try:
        file_size = file_path.stat().st_size
    except OSError as exc:
        raise DocumentProcessingError(
            f"Cannot read file: {exc}",
            filename=filename,
            stage="loader",
        ) from exc

    extra_meta: dict[str, str] = {}

    try:
        if suffix == ".pdf":
            text, page_count, extra_meta = await _load_pdf(file_path)
        elif suffix == ".docx":
            text, page_count = await _load_docx(file_path)
        elif suffix == ".xlsx":
            text, page_count = await _load_xlsx(file_path)
        elif suffix == ".pptx":
            text, page_count = await _load_pptx(file_path)
        elif suffix == ".csv":
            text, page_count = await _load_csv(file_path)
        else:
            raise UnsupportedFormatError(suffix, filename=filename)
    except (UnsupportedFormatError, DocumentProcessingError):
        raise
    except Exception as exc:
        raise DocumentProcessingError(
            f"Failed to load {suffix} file: {exc}",
            filename=filename,
            stage="loader",
        ) from exc

    word_count = _count_words(text)

    metadata = DocumentMetadata(
        filename=filename,
        file_type=suffix.lstrip("."),
        file_size_bytes=file_size,
        page_count=page_count,
        word_count=word_count,
        extra=extra_meta,
    )

    logger.info(
        "Loaded document '%s': %s, %d bytes, %s pages, %d words",
        filename,
        suffix,
        file_size,
        page_count or "N/A",
        word_count,
    )

    return text, metadata


async def _load_pdf(file_path: Path) -> tuple[str, int | None, dict[str, str]]:
    """Load a PDF file with tiered OCR fallback.

    Tier 1: PyPDFLoader (text layer extraction — always available).
    Tier 2: pdfplumber (better table/layout extraction — optional).
    Tier 3: PyMuPDF OCR (image-based OCR — optional, needs Tesseract).

    Returns the best text found across all attempted tiers plus extraction
    metadata indicating which method produced the result.
    """
    settings = get_settings()
    threshold = settings.pdf_ocr_min_text_chars
    ocr_enabled = settings.pdf_ocr_enabled

    tiers_attempted: list[str] = []
    best_text = ""
    best_method = "pypdf"
    page_count: int | None = None

    # ── Tier 1: PyPDFLoader (always runs) ──
    try:
        from langchain_community.document_loaders import PyPDFLoader

        loader = PyPDFLoader(str(file_path))
        pages: list[LCDocument] = await loader.aload()
        page_count = len(pages) if pages else None

        text_parts = [page.page_content for page in pages if page.page_content.strip()]
        tier1_text = "\f\n\n".join(text_parts)
        tiers_attempted.append("pypdf")

        if len(tier1_text.strip()) >= threshold:
            extra_meta = {
                "pdf_extraction_method": "pypdf",
                "pdf_ocr_applied": "false",
                "pdf_tiers_attempted": ",".join(tiers_attempted),
            }
            return tier1_text, page_count, extra_meta

        best_text = tier1_text
        logger.info(
            "PDF Tier 1 (pypdf) extracted %d chars (threshold %d) — trying fallback",
            len(tier1_text.strip()),
            threshold,
        )
    except Exception:
        tiers_attempted.append("pypdf")
        logger.warning("PDF Tier 1 (pypdf) failed", exc_info=True)

    if not ocr_enabled:
        extra_meta = {
            "pdf_extraction_method": best_method,
            "pdf_ocr_applied": "false",
            "pdf_tiers_attempted": ",".join(tiers_attempted),
        }
        return best_text, page_count, extra_meta

    # ── Tier 2: pdfplumber (optional) ──
    if _pdfplumber_available():
        try:
            tier2_text = await asyncio.to_thread(_extract_with_pdfplumber_sync, file_path)
            tiers_attempted.append("pdfplumber")

            if len(tier2_text.strip()) > len(best_text.strip()):
                best_text = tier2_text
                best_method = "pdfplumber"

            if len(best_text.strip()) >= threshold:
                extra_meta = {
                    "pdf_extraction_method": "pdfplumber",
                    "pdf_ocr_applied": "false",
                    "pdf_tiers_attempted": ",".join(tiers_attempted),
                }
                return best_text, page_count, extra_meta

            logger.info(
                "PDF Tier 2 (pdfplumber) extracted %d chars — trying OCR fallback",
                len(tier2_text.strip()),
            )
        except Exception:
            tiers_attempted.append("pdfplumber")
            logger.warning("PDF Tier 2 (pdfplumber) failed", exc_info=True)
    else:
        logger.debug("PDF Tier 2 skipped: pdfplumber not installed")

    # ── Tier 3: PyMuPDF OCR (optional) ──
    if _pymupdf_available():
        try:
            tier3_text = await asyncio.to_thread(_extract_with_pymupdf_ocr_sync, file_path)
            tiers_attempted.append("pymupdf_ocr")

            if len(tier3_text.strip()) > len(best_text.strip()):
                best_text = tier3_text
                best_method = "pymupdf_ocr"

            extra_meta = {
                "pdf_extraction_method": best_method,
                "pdf_ocr_applied": "true",
                "pdf_tiers_attempted": ",".join(tiers_attempted),
            }
            return best_text, page_count, extra_meta
        except Exception:
            tiers_attempted.append("pymupdf_ocr")
            logger.warning("PDF Tier 3 (PyMuPDF OCR) failed", exc_info=True)
    else:
        logger.debug("PDF Tier 3 skipped: PyMuPDF not installed")

    # ── All tiers exhausted — return best result ──
    extra_meta = {
        "pdf_extraction_method": best_method,
        "pdf_ocr_applied": str(best_method == "pymupdf_ocr").lower(),
        "pdf_tiers_attempted": ",".join(tiers_attempted),
    }
    return best_text, page_count, extra_meta


async def _load_docx(file_path: Path) -> tuple[str, int | None]:
    """Load a DOCX file using docx2txt."""
    import docx2txt

    text = docx2txt.process(str(file_path))
    return text or "", None


async def _load_xlsx(file_path: Path) -> tuple[str, int | None]:
    """Load an XLSX file using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    text_parts: list[str] = []
    sheet_count = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_count += 1
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(cell) if cell is not None else "" for cell in row]
            if any(cell_values):
                rows.append(" | ".join(cell_values))
        if rows:
            text_parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(text_parts), sheet_count if sheet_count else None


async def _load_pptx(file_path: Path) -> tuple[str, int | None]:
    """Load a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    text_parts: list[str] = []
    slide_count = 0

    for slide in prs.slides:
        slide_count += 1
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            text_parts.append(f"## Slide {slide_count}\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts), slide_count if slide_count else None


async def _load_csv(file_path: Path) -> tuple[str, int | None]:
    """Load a CSV file and convert to readable text."""
    with open(file_path, encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows: list[str] = []
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(row))

    return "\n".join(rows), None
